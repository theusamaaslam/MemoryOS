from __future__ import annotations

import csv
import hashlib
import json
import re
from dataclasses import dataclass, field
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text


LEGACY_OFFICE_EXTENSIONS = {".doc", ".ppt", ".xls"}
TEXT_EXTENSIONS = {
    ".txt",
    ".text",
    ".md",
    ".markdown",
    ".json",
    ".csv",
    ".html",
    ".htm",
    ".xml",
    ".yml",
    ".yaml",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".css",
    ".sql",
    ".sh",
    ".ini",
    ".cfg",
    ".toml",
}


class UnsupportedDocumentError(ValueError):
    pass


class DocumentParsingError(ValueError):
    pass


@dataclass(slots=True)
class ExtractedBlock:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class ParsedDocument:
    source_name: str
    source_uri: str
    source_type: str
    parser_name: str
    title: str | None
    blocks: list[ExtractedBlock]
    metadata: dict[str, Any] = field(default_factory=dict)


class DocumentIngestionService:
    def build_manual_document(self, source_name: str, text: str) -> ParsedDocument:
        cleaned = self._normalize_text(text)
        if not cleaned:
            raise DocumentParsingError("No usable content was found in the submitted text.")
        resolved_name = source_name.strip() or "Manual knowledge drop"
        return ParsedDocument(
            source_name=resolved_name,
            source_uri=f"manual://{self._slugify(resolved_name)}",
            source_type="manual_text",
            parser_name="plain_text",
            title=resolved_name,
            blocks=[ExtractedBlock(text=cleaned, metadata={"block_type": "manual_text"})],
        )

    def parse_upload(self, filename: str | None, content_type: str | None, payload: bytes) -> ParsedDocument:
        safe_name = Path(filename or "uploaded-document").name or "uploaded-document"
        extension = Path(safe_name).suffix.lower()
        source_uri = f"upload://{safe_name}"
        file_metadata = {
            "content_type": content_type or "application/octet-stream",
            "extension": extension or "",
            "file_size_bytes": len(payload),
        }

        if extension in LEGACY_OFFICE_EXTENSIONS:
            raise UnsupportedDocumentError(
                f"{extension} files are not supported yet. Export the document as PDF, DOCX, PPTX, or XLSX and try again."
            )
        if extension == ".pdf":
            return self._parse_pdf(safe_name, source_uri, payload, file_metadata)
        if extension in {".docx", ".docm"}:
            return self._parse_docx(safe_name, source_uri, payload, file_metadata)
        if extension == ".pptx":
            return self._parse_pptx(safe_name, source_uri, payload, file_metadata)
        if extension == ".xlsx":
            return self._parse_xlsx(safe_name, source_uri, payload, file_metadata)
        if extension == ".rtf":
            return self._parse_rtf(safe_name, source_uri, payload, file_metadata)
        if extension in {".html", ".htm"}:
            return self._parse_html(safe_name, source_uri, payload, file_metadata)
        if extension == ".xml":
            return self._parse_xml(safe_name, source_uri, payload, file_metadata)
        if extension == ".json":
            return self._parse_json(safe_name, source_uri, payload, file_metadata)
        if extension == ".csv":
            return self._parse_csv(safe_name, source_uri, payload, file_metadata)
        if extension in TEXT_EXTENSIONS or (content_type or "").startswith("text/") or not extension:
            return self._parse_text(safe_name, source_uri, payload, file_metadata)
        raise UnsupportedDocumentError(
            f"{safe_name} is not a supported upload type yet. Try PDF, DOCX, PPTX, XLSX, RTF, HTML, JSON, CSV, XML, or a text-like file."
        )

    def chunk_document(self, document: ParsedDocument) -> tuple[list[dict[str, Any]], dict[str, int | str]]:
        normalized_blocks = []
        for block in document.blocks:
            cleaned = self._normalize_text(block.text)
            if cleaned:
                normalized_blocks.append(ExtractedBlock(text=cleaned, metadata=block.metadata))
        if not normalized_blocks:
            raise DocumentParsingError("The uploaded file did not contain any extractable text.")

        total_characters = sum(len(block.text) for block in normalized_blocks)
        target_size, overlap_size = self._resolve_chunk_config(total_characters)

        segments: list[tuple[str, dict[str, Any]]] = []
        for block_index, block in enumerate(normalized_blocks, start=1):
            for segment_index, segment_text in enumerate(self._split_block(block.text, target_size), start=1):
                metadata = dict(block.metadata)
                metadata["block_index"] = block_index
                metadata["segment_index"] = segment_index
                segments.append((segment_text, metadata))

        chunks: list[dict[str, Any]] = []
        current_parts: list[str] = []
        current_metadata: list[dict[str, Any]] = []

        for segment_text, segment_metadata in segments:
            candidate = "\n\n".join([*current_parts, segment_text]).strip()
            if current_parts and len(candidate) > target_size:
                chunk_text = "\n\n".join(current_parts).strip()
                if chunk_text:
                    chunks.append(
                        self._build_chunk(
                            document=document,
                            chunk_index=len(chunks) + 1,
                            content=chunk_text,
                            segment_metadata=current_metadata,
                            target_size=target_size,
                            overlap_size=overlap_size,
                            total_characters=total_characters,
                        )
                    )
                overlap_seed = self._extract_overlap(chunk_text, overlap_size)
                current_parts = [overlap_seed] if overlap_seed else []
                current_metadata = []

            current_parts.append(segment_text)
            current_metadata.append(segment_metadata)

        if current_metadata:
            chunk_text = "\n\n".join(current_parts).strip()
            if chunk_text:
                chunks.append(
                    self._build_chunk(
                        document=document,
                        chunk_index=len(chunks) + 1,
                        content=chunk_text,
                        segment_metadata=current_metadata,
                        target_size=target_size,
                        overlap_size=overlap_size,
                        total_characters=total_characters,
                    )
                )

        for index, chunk in enumerate(chunks, start=1):
            chunk["metadata"]["chunk_index"] = index
            chunk["metadata"]["chunk_count"] = len(chunks)

        strategy = {
            "chunking_strategy": "dynamic_v1",
            "target_size": target_size,
            "overlap_size": overlap_size,
            "total_characters": total_characters,
        }
        return chunks, strategy

    def build_document_fingerprint(
        self,
        source_name: str,
        source_uri: str,
        source_type: str,
        chunks: list[dict[str, Any]],
        *,
        metadata: dict[str, Any] | None = None,
    ) -> str:
        payload = {
            "source_name": str(source_name or "").strip(),
            "source_uri": str(source_uri or "").strip(),
            "source_type": str(source_type or "").strip(),
            "metadata": self._stable_metadata(metadata or {}),
            "chunks": [
                {
                    "content": self._normalize_text(str(chunk.get("content", ""))),
                    "source_uri": str(chunk.get("source_uri", "")).strip(),
                    "title": str(chunk.get("title", "") or "").strip(),
                    "metadata": self._stable_metadata(chunk.get("metadata", {}) or {}),
                }
                for chunk in chunks
            ],
        }
        return hashlib.sha1(json.dumps(payload, sort_keys=True, ensure_ascii=True).encode("utf-8")).hexdigest()

    def build_chunk_key(
        self,
        *,
        source_id: str,
        source_uri: str,
        content: str,
        metadata: dict[str, Any] | None = None,
    ) -> str:
        payload = {
            "source_id": source_id,
            "source_uri": str(source_uri or "").strip(),
            "title_hint": self._normalize_text(str((metadata or {}).get("title") or "")),
            "metadata": {
                key: value
                for key, value in self._stable_metadata(metadata or {}).items()
                if key
                in {
                    "chunk_index",
                    "chunk_count",
                    "section_heading",
                    "section_headings",
                    "page_number",
                    "page_numbers",
                    "slide_number",
                    "slide_numbers",
                    "sheet_name",
                    "sheet_names",
                    "block_type",
                    "block_types",
                    "table_index",
                    "table_indices",
                    "structural_path",
                }
            },
        }
        digest = hashlib.sha1(json.dumps(payload, sort_keys=True, ensure_ascii=True).encode("utf-8")).hexdigest()
        return f"chunk-{digest[:40]}"

    def _parse_pdf(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        try:
            reader = PdfReader(BytesIO(payload))
        except Exception as exc:  # pragma: no cover - parser library exceptions vary
            raise DocumentParsingError(f"Could not parse PDF {source_name}.") from exc

        blocks = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = self._normalize_text(page.extract_text() or "")
            if text:
                blocks.append(ExtractedBlock(text=text, metadata={"block_type": "page", "page_number": page_number}))

        title = None
        if reader.metadata is not None:
            raw_title = getattr(reader.metadata, "title", None) or reader.metadata.get("/Title")
            if raw_title:
                title = str(raw_title).strip()

        return ParsedDocument(
            source_name=source_name,
            source_uri=source_uri,
            source_type="pdf",
            parser_name="pypdf",
            title=title or Path(source_name).stem,
            blocks=blocks,
            metadata=file_metadata,
        )

    def _parse_docx(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        try:
            document = DocxDocument(BytesIO(payload))
        except Exception as exc:  # pragma: no cover - parser library exceptions vary
            raise DocumentParsingError(f"Could not parse Word document {source_name}.") from exc

        blocks: list[ExtractedBlock] = []
        active_heading: str | None = None

        for paragraph in document.paragraphs:
            text = self._normalize_text(paragraph.text or "")
            if not text:
                continue
            style_name = (paragraph.style.name if paragraph.style else "").lower()
            if "heading" in style_name or style_name.startswith("title"):
                active_heading = text
                blocks.append(ExtractedBlock(text=text, metadata={"block_type": "heading", "section_heading": text}))
                continue

            metadata: dict[str, Any] = {"block_type": "paragraph"}
            if active_heading:
                metadata["section_heading"] = active_heading
            blocks.append(ExtractedBlock(text=text, metadata=metadata))

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                values = [self._normalize_text(cell.text or "") for cell in row.cells]
                values = [value for value in values if value]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                metadata = {"block_type": "table", "table_index": table_index}
                if active_heading:
                    metadata["section_heading"] = active_heading
                blocks.append(ExtractedBlock(text="\n".join(rows), metadata=metadata))

        title = self._first_non_empty([document.core_properties.title, Path(source_name).stem])
        return ParsedDocument(
            source_name=source_name,
            source_uri=source_uri,
            source_type="docx",
            parser_name="python-docx",
            title=title,
            blocks=blocks,
            metadata=file_metadata,
        )

    def _parse_pptx(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        try:
            presentation = Presentation(BytesIO(payload))
        except Exception as exc:  # pragma: no cover - parser library exceptions vary
            raise DocumentParsingError(f"Could not parse PowerPoint {source_name}.") from exc

        blocks: list[ExtractedBlock] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = []
            slide_title = None
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is not None:
                slide_title = self._normalize_text(getattr(title_shape, "text", "") or "")
            for shape in slide.shapes:
                raw_text = getattr(shape, "text", None)
                if raw_text:
                    cleaned = self._normalize_text(str(raw_text))
                    if cleaned:
                        texts.append(cleaned)
            if texts:
                metadata: dict[str, Any] = {"block_type": "slide", "slide_number": slide_number}
                if slide_title:
                    metadata["section_heading"] = slide_title
                blocks.append(ExtractedBlock(text="\n\n".join(texts), metadata=metadata))

        title = self._first_non_empty([presentation.core_properties.title, Path(source_name).stem])
        return ParsedDocument(
            source_name=source_name,
            source_uri=source_uri,
            source_type="pptx",
            parser_name="python-pptx",
            title=title,
            blocks=blocks,
            metadata=file_metadata,
        )

    def _parse_xlsx(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        try:
            workbook = load_workbook(filename=BytesIO(payload), read_only=True, data_only=True)
        except Exception as exc:  # pragma: no cover - parser library exceptions vary
            raise DocumentParsingError(f"Could not parse spreadsheet {source_name}.") from exc

        blocks: list[ExtractedBlock] = []
        try:
            for worksheet in workbook.worksheets:
                rows = []
                row_count = 0
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value not in (None, "")]
                    if values:
                        row_count += 1
                        rows.append(" | ".join(values))
                if rows:
                    blocks.append(
                        ExtractedBlock(
                            text=f"Sheet: {worksheet.title}\n" + "\n".join(rows),
                            metadata={
                                "block_type": "sheet",
                                "sheet_name": worksheet.title,
                                "row_count": row_count,
                            },
                        )
                    )
        finally:
            workbook.close()

        return ParsedDocument(
            source_name=source_name,
            source_uri=source_uri,
            source_type="xlsx",
            parser_name="openpyxl",
            title=Path(source_name).stem,
            blocks=blocks,
            metadata=file_metadata,
        )

    def _parse_rtf(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        try:
            text = rtf_to_text(self._decode_text(payload))
        except Exception as exc:  # pragma: no cover - parser library exceptions vary
            raise DocumentParsingError(f"Could not parse RTF file {source_name}.") from exc
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type="rtf",
            parser_name="striprtf",
            text=text,
            metadata=file_metadata | {"block_type": "rtf"},
        )

    def _parse_html(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        text = self._decode_text(payload)
        soup = BeautifulSoup(text, "html.parser")
        title = self._normalize_text(soup.title.string if soup.title and soup.title.string else "") or Path(source_name).stem
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type="html",
            parser_name="beautifulsoup4",
            text=soup.get_text("\n"),
            metadata=file_metadata | {"block_type": "html", "section_heading": title},
            title=title,
        )

    def _parse_xml(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        text = self._decode_text(payload)
        try:
            root = ElementTree.fromstring(text)
            extracted = "\n".join(segment.strip() for segment in root.itertext() if segment and segment.strip())
        except ElementTree.ParseError:
            extracted = text
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type="xml",
            parser_name="xml.etree",
            text=extracted,
            metadata=file_metadata | {"block_type": "xml"},
        )

    def _parse_json(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        text = self._decode_text(payload)
        try:
            document = json.loads(text)
            extracted = json.dumps(document, indent=2, ensure_ascii=True)
        except json.JSONDecodeError:
            extracted = text
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type="json",
            parser_name="json",
            text=extracted,
            metadata=file_metadata | {"block_type": "json"},
        )

    def _parse_csv(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        text = self._decode_text(payload)
        reader = csv.reader(StringIO(text))
        rows = []
        row_count = 0
        for row in reader:
            values = [value.strip() for value in row if value and value.strip()]
            if values:
                row_count += 1
                rows.append(" | ".join(values))
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type="csv",
            parser_name="csv",
            text="\n".join(rows),
            metadata=file_metadata | {"block_type": "csv", "row_count": row_count},
        )

    def _parse_text(
        self,
        source_name: str,
        source_uri: str,
        payload: bytes,
        file_metadata: dict[str, Any],
    ) -> ParsedDocument:
        extension = Path(source_name).suffix.lower().lstrip(".") or "text"
        return self._single_block_document(
            source_name=source_name,
            source_uri=source_uri,
            source_type=extension,
            parser_name="plain_text",
            text=self._decode_text(payload),
            metadata=file_metadata | {"block_type": "text"},
        )

    def _single_block_document(
        self,
        source_name: str,
        source_uri: str,
        source_type: str,
        parser_name: str,
        text: str,
        metadata: dict[str, Any],
        title: str | None = None,
    ) -> ParsedDocument:
        cleaned = self._normalize_text(text)
        if not cleaned:
            raise DocumentParsingError(f"{source_name} did not contain any extractable text.")
        return ParsedDocument(
            source_name=source_name,
            source_uri=source_uri,
            source_type=source_type,
            parser_name=parser_name,
            title=title or Path(source_name).stem,
            blocks=[ExtractedBlock(text=cleaned, metadata=metadata)],
            metadata=metadata,
        )

    def _build_chunk(
        self,
        document: ParsedDocument,
        chunk_index: int,
        content: str,
        segment_metadata: list[dict[str, Any]],
        target_size: int,
        overlap_size: int,
        total_characters: int,
    ) -> dict[str, Any]:
        metadata: dict[str, Any] = {
            "kind": "document_chunk",
            "chunking_strategy": "dynamic_v1",
            "parser": document.parser_name,
            "source_type": document.source_type,
            "source_name": document.source_name,
            "target_chunk_size": target_size,
            "overlap_size": overlap_size,
            "document_characters": total_characters,
            "segment_count": len(segment_metadata),
        }
        metadata.update(document.metadata)

        for source_key, target_key in (
            ("page_number", "page_numbers"),
            ("slide_number", "slide_numbers"),
            ("sheet_name", "sheet_names"),
            ("section_heading", "section_headings"),
            ("block_type", "block_types"),
            ("table_index", "table_indices"),
        ):
            values = self._unique_values(meta.get(source_key) for meta in segment_metadata)
            if values:
                metadata[target_key] = values
        structural_parts = []
        for key in ("section_headings", "sheet_names", "page_numbers", "slide_numbers", "block_types"):
            value = metadata.get(key)
            if isinstance(value, list) and value:
                structural_parts.append(f"{key}:{','.join(str(part) for part in value[:4])}")
        if structural_parts:
            metadata["structural_path"] = " | ".join(structural_parts)

        return {
            "content": content,
            "source_uri": f"{document.source_uri}#chunk={chunk_index}",
            "title": document.title or document.source_name,
            "metadata": metadata,
        }

    def _resolve_chunk_config(self, total_characters: int) -> tuple[int, int]:
        if total_characters <= 3_000:
            return 700, 100
        if total_characters <= 12_000:
            return 1_000, 140
        if total_characters <= 40_000:
            return 1_300, 180
        return 1_700, 220

    def _split_block(self, text: str, target_size: int) -> list[str]:
        paragraphs = [segment.strip() for segment in re.split(r"\n{2,}", text) if segment.strip()]
        if not paragraphs:
            return []

        segments: list[str] = []
        for paragraph in paragraphs:
            if len(paragraph) <= target_size:
                segments.append(paragraph)
                continue
            segments.extend(self._split_large_paragraph(paragraph, target_size))
        return segments

    def _split_large_paragraph(self, paragraph: str, target_size: int) -> list[str]:
        sentences = [segment.strip() for segment in re.split(r"(?<=[.!?])\s+", paragraph) if segment.strip()]
        if len(sentences) <= 1:
            return self._hard_split(paragraph, target_size)

        segments: list[str] = []
        buffer: list[str] = []
        for sentence in sentences:
            if len(sentence) > target_size:
                if buffer:
                    segments.append(" ".join(buffer).strip())
                    buffer = []
                segments.extend(self._hard_split(sentence, target_size))
                continue
            candidate = " ".join([*buffer, sentence]).strip()
            if buffer and len(candidate) > target_size:
                segments.append(" ".join(buffer).strip())
                buffer = [sentence]
            else:
                buffer.append(sentence)
        if buffer:
            segments.append(" ".join(buffer).strip())
        return segments

    def _hard_split(self, text: str, target_size: int) -> list[str]:
        words = text.split()
        if not words:
            return []

        segments: list[str] = []
        buffer: list[str] = []
        for word in words:
            candidate = " ".join([*buffer, word]).strip()
            if buffer and len(candidate) > target_size:
                segments.append(" ".join(buffer).strip())
                buffer = [word]
            else:
                buffer.append(word)
        if buffer:
            segments.append(" ".join(buffer).strip())
        return segments

    def _extract_overlap(self, text: str, overlap_size: int) -> str:
        if not text:
            return ""
        if len(text) <= overlap_size:
            return text
        tail = text[-overlap_size:]
        sentence_start = tail.find(". ")
        paragraph_start = tail.find("\n\n")
        candidates = [index for index in (sentence_start, paragraph_start) if index >= 0]
        if candidates:
            boundary = min(candidates)
            return tail[boundary + 2 :].strip()
        return tail.strip()

    def _decode_text(self, payload: bytes) -> str:
        for encoding in ("utf-8", "utf-8-sig", "utf-16", "utf-16-le", "utf-16-be", "latin-1"):
            try:
                return payload.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise DocumentParsingError("Could not decode the uploaded file as text.")

    def _normalize_text(self, text: str) -> str:
        text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\u00a0", " ").replace("\x00", " ")
        lines = []
        blank_run = 0
        for raw_line in text.split("\n"):
            cleaned = re.sub(r"[ \t]+", " ", raw_line).strip()
            if not cleaned:
                blank_run += 1
                if blank_run <= 2:
                    lines.append("")
                continue
            blank_run = 0
            lines.append(cleaned)
        normalized = "\n".join(lines).strip()
        return re.sub(r"\n{3,}", "\n\n", normalized)

    def _slugify(self, value: str) -> str:
        normalized = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
        return normalized or "document"

    def _unique_values(self, values: Any) -> list[Any]:
        result = []
        for value in values:
            if value in (None, "", []):
                continue
            if value not in result:
                result.append(value)
        return result

    def _first_non_empty(self, values: list[Any]) -> str | None:
        for value in values:
            if value:
                resolved = str(value).strip()
                if resolved:
                    return resolved
        return None

    def _stable_metadata(self, metadata: dict[str, Any]) -> dict[str, Any]:
        stable: dict[str, Any] = {}
        for key, value in (metadata or {}).items():
            if isinstance(value, dict):
                stable[str(key)] = self._stable_metadata(value)
            elif isinstance(value, list):
                stable[str(key)] = [
                    self._stable_metadata(item)
                    if isinstance(item, dict)
                    else item
                    if isinstance(item, (str, int, float, bool)) or item is None
                    else str(item)
                    for item in value
                ]
            elif isinstance(value, (str, int, float, bool)) or value is None:
                stable[str(key)] = value
            else:
                stable[str(key)] = str(value)
        return stable


document_ingestion_service = DocumentIngestionService()
